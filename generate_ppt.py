import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Industrial Precision Dark Theme
    COLOR_BG = RGBColor(10, 14, 20)          # #0a0e14
    COLOR_SURFACE = RGBColor(18, 25, 36)     # #121924
    COLOR_BORDER = RGBColor(30, 41, 59)      # #1e293b
    COLOR_ACCENT = RGBColor(184, 255, 74)    # #b8ff4a (ProdNexus Green)
    COLOR_CYAN = RGBColor(56, 189, 248)      # #38bdf8 (Tech Cyan)
    COLOR_WHITE = RGBColor(248, 250, 252)    # #f8fafc
    COLOR_MUTED = RGBColor(148, 163, 184)    # #94a3b8

    blank_layout = prs.slide_layouts[6]

    def base_slide(title_text="", category_text=""):
        slide = prs.slides.add_slide(blank_layout)
        
        # Solid Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Minimal Top Indicator Bar
        top_accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
        top_accent.fill.solid()
        top_accent.fill.fore_color.rgb = COLOR_ACCENT
        top_accent.line.fill.background()

        if title_text:
            tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
            tf = tb.text_frame
            tf.word_wrap = True
            
            p_cat = tf.paragraphs[0]
            p_cat.text = (category_text or "UNIHACK 2026 · AI-POWERED PRODUCT INTELLIGENCE").upper()
            p_cat.font.size = Pt(9.5)
            p_cat.font.bold = True
            p_cat.font.color.rgb = COLOR_ACCENT

            p_title = tf.add_paragraph()
            p_title.text = title_text
            p_title.font.size = Pt(21)
            p_title.font.bold = True
            p_title.font.color.rgb = COLOR_WHITE
            p_title.space_before = Pt(3)

        return slide

    def card(slide, left, top, width, height, title, items, border_color=COLOR_BORDER, title_color=COLOR_CYAN):
        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_SURFACE
        c.line.color.rgb = border_color
        c.line.width = Pt(1)

        tb = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.18), width - Inches(0.44), height - Inches(0.36))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = title.upper()
        p_t.font.size = Pt(12.5)
        p_t.font.bold = True
        p_t.font.color.rgb = title_color

        for item in items:
            p = tf.add_paragraph()
            if isinstance(item, tuple):
                p.text = f"{item[0]}: "
                p.font.bold = True
                p.font.size = Pt(10.5)
                p.font.color.rgb = COLOR_WHITE
                p.space_before = Pt(5)
                
                # Add remainder of text
                run = p.add_run()
                run.text = item[1]
                run.font.bold = False
                run.font.color.rgb = COLOR_MUTED
            else:
                p.text = f"• {item}"
                p.font.size = Pt(10.5)
                p.font.color.rgb = COLOR_MUTED
                p.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 1: Title & Team Details
    # -------------------------------------------------------------
    s1 = base_slide()
    tb1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(3.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "UNIHACK 2026 PROTOTYPE SUBMISSION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    p = tf1.add_paragraph()
    p.text = "ProdNexus"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(6)

    p = tf1.add_paragraph()
    p.text = "Autonomous Industrial Product Intelligence & Decision-Support System"
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_CYAN
    p.space_before = Pt(8)

    card(s1, Inches(0.8), Inches(4.5), Inches(5.6), Inches(2.2), "TEAM DETAILS", [
        ("Team Name", "Team Catalyst"),
        ("Team Leader", "Muskan Mulani"),
        ("Track", "AI-Powered Product Intelligence for Industrial Commerce"),
        ("Target Users", "Procurement Engineers, Catalog Managers, Industrial Buyers")
    ], COLOR_ACCENT, COLOR_ACCENT)

    card(s1, Inches(6.8), Inches(4.5), Inches(5.7), Inches(2.2), "CORE TECHNICAL HIGHLIGHTS", [
        ("Sub-3s Inference", "Deterministic latency via Gemini 2.5 Flash + pruned payloads"),
        ("AI Decision Score", "Automated 0–100 index calculating fit %, spec advantage & supply risk"),
        ("Enterprise Ingestion", "Handles both single SKU deep-dives and 50+ row CSV batch queues"),
        ("Interactive Explainer", "Context-grounded slide-over chat drawer for real-time auditability")
    ], COLOR_BORDER, COLOR_CYAN)

    # -------------------------------------------------------------
    # SLIDE 2: Problem vs. Solution
    # -------------------------------------------------------------
    s2 = base_slide("Problem Definition & Solution Overview", "01 / EXECUTIVE CONTEXT")
    
    card(s2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), "THE INDUSTRIAL COMMERCE BOTTLENECK", [
        ("Fragmented Catalog Data", "Industrial SKUs often arrive with only an MPN and a 4-word snippet, missing standard voltage, mounting, or IP ratings."),
        ("Manual Datasheet Cross-Referencing", "Engineers spend 15–30 minutes per SKU verifying OEM alternatives across disparate PDF catalogs."),
        ("Vendor Pricing Opacity", "Over 60% of industrial components lack public list prices, halting procurement workflows for quote requests."),
        ("Lack of Decision Frameworks", "Existing PIMs act as passive databases; they don't explain *which* product to select for a given requirement.")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s2, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), "THE PRODNEXUS VALUE PROPOSITION", [
        ("Sub-Second Semantic Retrieval", "Cosine-similarity vector retrieval extracts real catalog equivalents instantaneously."),
        ("Structured AI Attribute Extraction", "Enforces strict JSON schema to generate normalized summaries, features, strengths, and limits."),
        ("Benchmark Price Range Estimator", "AI market bracket estimator provides realistic enterprise procurement estimates tagged explicitly."),
        ("Quantified Decision Index", "Scores candidates from 0–100 with automated risk verdicts (STRONG CANDIDATE vs NEEDS REVIEW).")
    ], COLOR_ACCENT, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 3: Technical Responses to Guidelines
    # -------------------------------------------------------------
    s3 = base_slide("Technical Core & Enterprise Architecture", "02 / EVALUATION QUESTIONS")

    card(s3, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1), "1. ENRICHMENT STRATEGY", [
        ("Sparse Input Ingestion", "Accepts minimal MPN, Brand, and Short Description."),
        ("Vector Context Assembly", "Retrieves top-k similar catalog products and trims metadata to 4 lightweight matches."),
        ("Constrained LLM Reasoning", "Prompts Gemini 2.5 Flash to extract operational voltage, IP ratings, and duty classes."),
        ("Structured JSON Output", "Enforces deterministic JSON schemas without conversational tokens.")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s3, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1), "2. ACCURACY & TRUST", [
        ("RAG Grounding", "Specs are strictly bound to retrieved catalog items to prevent hallucination."),
        ("Explicit Benchmark Tags", "Estimated pricing is clearly badged with ESTIMATED BENCHMARK tags."),
        ("Decision Metric Math", "Combines feature depth, strengths, and catalog coverage into a deterministic score."),
        ("Context-Aware Chat Drawer", "Allows procurement engineers to query reasons directly via Ask AI.")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s3, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.1), "3. ENTERPRISE SCALABILITY", [
        ("thinkingBudget: 0", "Disables chain-of-thought token bloat, reducing latency from 30s to <3.5s."),
        ("Asynchronous Batch Queue", "Processes 50+ row CSVs sequentially with visual progress tracking."),
        ("Failsafe Resilience", "Individual row errors trigger graceful fallbacks without crashing the batch."),
        ("Universal Normalization", "Normalizes irregular CSV header formats (Mfg_Part_Num, SKU, Part_Desc).")
    ], COLOR_BORDER, COLOR_CYAN)

    # -------------------------------------------------------------
    # SLIDE 4: Opportunities & USP
    # -------------------------------------------------------------
    s4 = base_slide("Opportunities & Unique Selling Proposition", "03 / MARKET ADVANTAGE")

    card(s4, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), "TRADITIONAL PIM VS. PRODNEXUS", [
        ("Traditional Search", "Matches exact keywords only; fails when part numbers differ slightly or abbreviations are used."),
        ("ProdNexus Semantic Search", "Understands technical synonyms, voltage equivalents, and compatible hardware classes."),
        ("Static Data Display", "Traditional systems display raw text with missing fields left blank."),
        ("Active Decision Engine", "ProdNexus benchmarks price brackets, scores supply risk, and generates procurement recommendations.")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s4, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), "KEY USPs (WHY PRODNEXUS WINS)", [
        ("⭐ Quantified Decision Score", "Generates an objective 0–100 score + verdict pill based on catalog metrics."),
        ("⚡ Production-Ready Latency", "End-to-end enrichment executes in 2–4 seconds per item."),
        ("📦 Dual Ingestion Workflows", "Interactive single-SKU analysis + Bulk CSV batch processor in one unified UI."),
        ("💬 Grounded Chat Explainer", "Slide-over drawer allows live verification of technical trade-offs."),
        ("📄 Executive Dossier Export", "Generates client-side formatted PDF procurement briefs on demand.")
    ], COLOR_ACCENT, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 5: List of Features
    # -------------------------------------------------------------
    s5 = base_slide("Key Functional Features", "04 / CAPABILITIES MATRIX")

    card(s5, Inches(0.8), Inches(1.6), Inches(3.7), Inches(2.4), "1. SINGLE SKU DEEP-DIVE", [
        ("Instant Profile", "Market tier & executive summaries"),
        ("Key Attributes", "Confirmed hardware specifications"),
        ("Pros & Cons", "Operational strengths & limitations")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s5, Inches(4.8), Inches(1.6), Inches(3.7), Inches(2.4), "2. BULK CSV PROCESSOR", [
        ("Dropzone Upload", "Auto-parses irregular industrial CSVs"),
        ("Queue Monitor", "Real-time percentage progress bar"),
        ("Export Tool", "Downloads enriched CSV with scores")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s5, Inches(8.8), Inches(1.6), Inches(3.7), Inches(2.4), "3. AI DECISION GAUGE", [
        ("Score: 0–100", "Metric-backed procurement score"),
        ("Fit Breakdown", "Market Fit % & Spec Advantage %"),
        ("Risk Badge", "LOW / MODERATE supply chain risk")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s5, Inches(0.8), Inches(4.3), Inches(3.7), Inches(2.4), "4. 3-WAY COMPARISON", [
        ("Vector Ranking", "Calculates % relevance to target SKU"),
        ("Comparison Shelf", "Pin up to 3 models side-by-side"),
        ("Brand Filters", "Dynamic brand chips & keyword filter")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s5, Inches(4.8), Inches(4.3), Inches(3.7), Inches(2.4), "5. 'ASK AI' EXPLAINER", [
        ("Slide-Over Drawer", "Embedded conversational assistant"),
        ("Context-Grounded", "Queries strictly against catalog data"),
        ("Audit Ready", "Explains selection rationales live")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s5, Inches(8.8), Inches(4.3), Inches(3.7), Inches(2.4), "6. PDF DOSSIER & HISTORY", [
        ("Offline Sharing", "Generates executive PDF reports"),
        ("Session Storage", "Local persistence of recent runs"),
        ("Row Inspect", "1-click inspection from batch tables")
    ], COLOR_BORDER, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 6: Process Flow Diagram
    # -------------------------------------------------------------
    s6 = base_slide("Operational Process Flow & Data Pipeline", "05 / SYSTEM WORKFLOW")

    card(s6, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.1), "END-TO-END PIPELINE ARCHITECTURE", [
        ("1. Input Ingestion", "User submits single SKU parameters or uploads a multi-row industrial catalog CSV file."),
        ("2. Schema Normalization", "Frontend / Controller normalizes column headers (Mfg_Part_Num, SKU, Part_Desc) into standard keys."),
        ("3. Semantic Vector Retrieval", "Retrieval service executes vector cosine matching against the catalog database to find relevant equivalents."),
        ("4. Prompt Context Pruning", "Retrieved matches are trimmed to top 4 items and essential fields to minimize token transmission latency."),
        ("5. Zero-Delay Inference", "Gemini 2.5 Flash generates structured intelligence (enforcing JSON schema with thinkingBudget: 0)."),
        ("6. Decision Metric Engine", "Backend calculates overall score (0–100), verdict badge, market fit, spec advantage, and supply risk."),
        ("7. UI Presentation & Export", "Client renders glowing score card, comparison shelf, Ask AI drawer, with 1-click CSV/PDF export.")
    ], COLOR_ACCENT, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 7: Technical Architecture
    # -------------------------------------------------------------
    s7 = base_slide("Technical Architecture Diagram", "06 / SYSTEM DESIGN")

    card(s7, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.1), "CLIENT LAYER (REACT 18 SPA)", [
        ("Component Framework", "React 18 with Vite fast tooling"),
        ("UI & Styling", "Custom modular dark design system, CSS Grid/Flexbox"),
        ("Document Generation", "jsPDF vector engine for client-side reports"),
        ("Session State", "LocalStorage API for persistent analysis history"),
        ("Batch Streaming", "Asynchronous state queues for real-time progress")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s7, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.1), "BACKEND & AI ORCHESTRATION LAYER", [
        ("Runtime & API Gateway", "Node.js & Express.js REST API"),
        ("AI Orchestrator", "Gemini 2.5 Flash via official @google/genai SDK"),
        ("Inference Controls", "temperature: 0.2, thinkingBudget: 0, strict JSON schema"),
        ("Retrieval Engine", "Vector similarity cosine index for sub-second SKU matching"),
        ("Decision Engine", "Deterministic heuristic scoring calculating fit % and risk tiers")
    ], COLOR_ACCENT, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 8: Technologies Used
    # -------------------------------------------------------------
    s8 = base_slide("Technologies & Libraries Used", "07 / TECH STACK")

    card(s8, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1), "FRONTEND STACK", [
        ("React.js (v18)", "Component-driven single page application"),
        ("Vite", "Next-generation frontend tooling & bundling"),
        ("CSS3 Modules", "Custom dark theme design system"),
        ("jsPDF", "Client-side procurement dossier export"),
        ("FileReader API", "Browser-side CSV streaming & parsing")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s8, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1), "BACKEND & DATA", [
        ("Node.js", "Asynchronous non-blocking backend runtime"),
        ("Express.js", "RESTful API routes & validation middleware"),
        ("Vector Indexing", "Catalog similarity scoring & ranking"),
        ("Batch Queueing", "Sequential batch enrichment pipeline"),
        ("CORS / Dotenv", "Secure environment configuration")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s8, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.1), "AI & DEPLOYMENT", [
        ("Gemini 2.5 Flash", "High-speed reasoning LLM"),
        ("@google/genai SDK", "Official Google AI integration"),
        ("JSON Schema Mode", "Guarantees valid parseable payloads"),
        ("Git & GitHub", "Version control & collaboration"),
        ("npm / Node", "Package management & environment")
    ], COLOR_BORDER, COLOR_CYAN)

    # -------------------------------------------------------------
    # SLIDE 9: Implementation Cost & ROI
    # -------------------------------------------------------------
    s9 = base_slide("Estimated Implementation & Operating Cost", "08 / FEASIBILITY & TCO")

    card(s9, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1), "MODEL INFERENCE COSTS", [
        ("Model Tier", "Gemini 2.5 Flash"),
        ("Pricing Benchmark", "~$0.075 / 1M input tokens"),
        ("Per-SKU Cost", "~$0.00015 per enrichment"),
        ("100,000 SKUs", "Less than $15.00 USD total AI compute cost")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s9, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1), "INFRASTRUCTURE & HOSTING", [
        ("Application Server", "Node.js container instance: ~$20 - $40/mo"),
        ("Frontend Hosting", "Edge CDN static deployment: ~$0 - $15/mo"),
        ("Vector Store Tier", "Serverless vector index: ~$0 - $60/mo"),
        ("Total Hosting", "Under $115.00 / month")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s9, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.1), "ENTERPRISE ROI IMPACT", [
        ("Labor Reduction", "Eliminates ~90% of manual datasheet lookup time"),
        ("Catalog Time-to-Market", "Enriches 10,000 SKUs in hours instead of weeks"),
        ("Procurement Speed", "Accelerates quote approval cycles by 4x"),
        ("Payback Period", "Full implementation ROI achieved in <14 days")
    ], COLOR_BORDER, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 10: Future Roadmap
    # -------------------------------------------------------------
    s10 = base_slide("Future Roadmap & Enterprise Scaling", "09 / PRODUCT ROADMAP")

    card(s10, Inches(0.8), Inches(1.6), Inches(3.7), Inches(5.1), "PHASE 1: ERP / PIM CONNECTORS", [
        ("Two-Way Sync", "Direct bidirectional sync with SAP S/4HANA & Oracle ERP"),
        ("PIM Export", "Automatic taxonomy mapping into Akeneo & Pimcore"),
        ("PO Triggers", "Auto-generates purchase requisitions for strong matches")
    ], COLOR_BORDER, COLOR_CYAN)

    card(s10, Inches(4.8), Inches(1.6), Inches(3.7), Inches(5.1), "PHASE 2: MULTIMODAL PARSING", [
        ("Datasheet Ingestion", "Upload raw PDF engineering datasheets directly"),
        ("Gemini Vision", "Extract wiring schematics & dimensional drawings"),
        ("CAD Verification", "Auto-validates mounting bolt patterns and specs")
    ], COLOR_BORDER, COLOR_ACCENT)

    card(s10, Inches(8.8), Inches(1.6), Inches(3.7), Inches(5.1), "PHASE 3: LIVE SUPPLY CHAIN", [
        ("Distributor APIs", "Live inventory feeds (Mouser, DigiKey, Grainger)"),
        ("Lead Time Alerts", "Real-time stock level factor into risk score"),
        ("Predictive Pricing", "Tracks historical component price volatility")
    ], COLOR_BORDER, COLOR_CYAN)

    # -------------------------------------------------------------
    # SLIDE 11: Submission Links
    # -------------------------------------------------------------
    s11 = base_slide("Project Links & Prototype Deliverables", "10 / SUBMISSION DELIVERABLES")

    card(s11, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.1), "PROJECT REPOSITORIES & DEMO ASSETS", [
        ("GitHub Public Repository", "https://github.com/Muskanmulani/ProdNexus"),
        ("Demo Video Link (3 Minutes)", "https://youtu.be/your-demo-link (Includes live Single SKU + Bulk CSV demo)"),
        ("Live Working Prototype", "http://localhost:5173 / Deployed Web Application"),
        ("Submission Track", "UniHack 2026 · AI-Powered Product Intelligence for Industrial Commerce"),
        ("Team Name", "Team Catalyst (Leader: Muskan Mulani)")
    ], COLOR_ACCENT, COLOR_ACCENT)

    # -------------------------------------------------------------
    # SLIDE 12: Thank You
    # -------------------------------------------------------------
    s12 = base_slide()
    tb12 = s12.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(3.5))
    tf12 = tb12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Thank You"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf12.add_paragraph()
    p.text = "ProdNexus — Autonomous Product Intelligence for Industrial Commerce"
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_CYAN
    p.space_before = Pt(8)

    p = tf12.add_paragraph()
    p.text = "Team Catalyst  |  Team Leader: Muskan Mulani  |  UniHack 2026"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_ACCENT
    p.space_before = Pt(14)

    output_filename = "ProdNexus_UniHack_Presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation generated: {output_filename}")

if __name__ == "__main__":
    build_presentation()